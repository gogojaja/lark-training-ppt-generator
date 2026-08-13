#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
技能包内容增强工具
为技能包添加有意义的模板、样式和示例
"""

import os
from pathlib import Path

# 技能包列表
SKILL_PACKAGES = [
    'cover-skill',
    'toc-skill',
    'scene-description-skill',
    'business-rules-skill',
    'operation-steps-skill',
    'key-points-skill',
    'faq-skill',
    'ppt-framework',
    'style-brief-skill',
    'table-skill',
    'document-processing',
    'flowchart-skill',
]

def enhance_cover_skill(skill_path):
    """增强封面技能包"""
    print(f"  [+] 增强封面技能包")

    # 更新 templates/template_example.txt
    template_file = skill_path / 'templates' / 'template_example.txt'
    template_content = """# 封面页模板

## 模板说明
标准封面页模板，包含标题、副标题、作者信息

## 使用方法
1. 替换 {{title}} 为演示主题
2. 替换 {{subtitle}} 为副标题
3. 替换 {{author}} 为作者/团队名称
4. 替换 {{date}} 为日期

## 示例
{{title}}: 个人综合签约培训
{{subtitle}}: 操作流程与注意事项
{{author}}: 培训部
{{date}}: 2026-08-13
"""
    template_file.write_text(template_content, encoding='utf-8')

    # 更新 styles/style_example.css
    style_file = skill_path / 'styles' / 'style_example.css'
    style_content = """/* 封面页样式定义 */

/* 颜色定义 */
--cover-primary: #1F3864;      /* 主色调 - 深蓝 */
--cover-secondary: #2E75B6;    /* 辅助色 - 中蓝 */
--cover-accent: #ED7D31;       /* 强调色 - 橙色 */
--cover-text: #333333;         /* 文本色 */
--cover-light: #F5F7FA;        /* 浅色背景 */
--cover-white: #FFFFFF;        /* 白色背景 */

/* 字体定义 */
--cover-title-font: 'Microsoft YaHei', '微软雅黑', sans-serif;
--cover-subtitle-font: 'Microsoft YaHei', '微软雅黑', sans-serif;
--cover-author-font: 'Arial', sans-serif;

/* 字号定义 */
--cover-title-size: 48px;      /* 标题字号 */
--cover-subtitle-size: 28px;   /* 副标题字号 */
--cover-author-size: 20px;     /* 作者字号 */

/* 间距定义 */
--cover-title-spacing: 20px;
--cover-subtitle-spacing: 10px;
--cover-author-spacing: 15px;

/* 布局定义 */
--cover-padding: 60px;
--cover-header-height: 80px;
--cover-footer-height: 60px;

/* 效果定义 */
--cover-blend-opacity: 0.88;   /* 混合透明度 */
--cover-shadow: 0 4px 12px rgba(0, 0, 0, 0.1);
--cover-gradient: linear-gradient(135deg, var(--cover-primary) 0%, var(--cover-secondary) 100%);
"""
    style_file.write_text(style_content, encoding='utf-8')

    # 更新 examples/example.py
    example_file = skill_path / 'examples' / 'example.py'
    example_content = """# 封面页示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_cover_page(pres, title, subtitle, author, date):
    \"\"\"创建封面页\"\"\"
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 设置背景色（深蓝渐变）
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 56, 100)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle

    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)

    # 作者
    author_box = slide.shapes.add_textbox(
        Inches(1), Inches(6), Inches(8), Inches(0.5)
    )
    author_frame = author_box.text_frame
    author_frame.text = author

    author_para = author_frame.paragraphs[0]
    author_para.alignment = PP_ALIGN.CENTER
    author_para.font.size = Pt(20)
    author_para.font.color.rgb = RGBColor(200, 200, 200)

    # 日期
    date_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5), Inches(8), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_frame.text = date

    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(20)
    date_para.font.color.rgb = RGBColor(200, 200, 200)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    slide = create_cover_page(
        pres,
        title='个人综合签约培训',
        subtitle='操作流程与注意事项',
        author='培训部',
        date='2026-08-13'
    )
    pres.save('output/cover_example.pptx')
    print('封面页示例已生成: output/cover_example.pptx')
"""
    example_file.write_text(example_content, encoding='utf-8')

def enhance_toc_skill(skill_path):
    """增强目录技能包"""
    print(f"  [+] 增强目录技能包")

    # 更新 templates/template_example.txt
    template_file = skill_path / 'templates' / 'template_example.txt'
    template_content = """# 目录页模板

## 模板说明
标准目录页模板，展示培训内容的结构

## 使用方法
1. 替换 {{title}} 为目录标题
2. 替换 {{items}} 为目录项列表
3. 可选：添加页码范围

## 示例
{{title}}: 目录
{{items}}:
1. 个人综合签约概述 - 第1-2页
2. 签约流程详解 - 第3-8页
3. 注意事项 - 第9-12页
4. 常见问题 - 第13-15页
"""
    template_file.write_text(template_content, encoding='utf-8')

    # 更新 styles/style_example.css
    style_file = skill_path / 'styles' / 'style_example.css'
    style_content = """/* 目录页样式定义 */

/* 颜色定义 */
--toc-primary: #1F3864;
--toc-secondary: #2E75B6;
--toc-text: #333333;
--toc-light: #F5F7FA;
--toc-border: #E0E0E0;

/* 字体定义 */
--toc-title-font: 'Microsoft YaHei', '微软雅黑', sans-serif;
--toc-item-font: 'Microsoft YaHei', '微软雅黑', sans-serif;

/* 字号定义 */
--toc-title-size: 36px;
--toc-item-size: 24px;
--toc-page-size: 18px;

/* 间距定义 */
--toc-title-spacing: 40px;
--toc-item-spacing: 12px;
--toc-item-padding: 8px 12px;

/* 布局定义 */
--toc-column-count: 2;
--toc-column-gap: 60px;
--toc-content-width: 600px;
"""
    style_file.write_text(style_content, encoding='utf-8')

    # 更新 examples/example.py
    example_file = skill_path / 'examples' / 'example.py'
    example_content = """# 目录页示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_toc_page(pres, title, items):
    \"\"\"创建目录页\"\"\"
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 56, 100)

    # 目录项
    start_y = Inches(2.5)
    x_position = Inches(2)

    for i, item in enumerate(items):
        # 创建文本框
        text_box = slide.shapes.add_textbox(
            x_position, start_y + i * Inches(0.6), Inches(6), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = item

        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(24)
        text_para.font.color.rgb = RGBColor(51, 51, 51)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    items = [
        '个人综合签约概述 - 第1-2页',
        '签约流程详解 - 第3-8页',
        '注意事项 - 第9-12页',
        '常见问题 - 第13-15页'
    ]
    slide = create_toc_page(pres, '目录', items)
    pres.save('output/toc_example.pptx')
    print('目录页示例已生成: output/toc_example.pptx')
"""
    example_file.write_text(example_content, encoding='utf-8')

def enhance_table_skill(skill_path):
    """增强表格技能包"""
    print(f"  [+] 增强表格技能包")

    # 更新 templates/template_example.txt
    template_file = skill_path / 'templates' / 'template_example.txt'
    template_content = """# 表格模板

## 模板说明
标准表格模板，展示数据对比和分类

## 使用方法
1. 替换 {{title}} 为表格标题
2. 替换 {{columns}} 为列名
3. 替换 {{rows}} 为数据行
4. 选择模板类型（标准/斑马纹/图标）

## 示例
{{title}}: 签约流程步骤对比
{{columns}}: 步骤 | 操作 | 时间 | 负责人
{{rows}}:
1 | 提交申请 | 1-2天 | 客户经理
2 | 审核材料 | 2-3天 | 风控专员
3 | 签订合同 | 1天 | 法务部门
"""
    template_file.write_text(template_content, encoding='utf-8')

    # 更新 styles/style_example.css
    style_file = skill_path / 'styles' / 'style_example.css'
    style_content = """/* 表格样式定义 */

/* 颜色定义 */
--table-primary: #1F3864;
--table-secondary: #2E75B6;
--table-header: #4A6FA5;
--table-border: #E0E0E0;
--table-bg: #FFFFFF;
--table-stripe: #F5F7FA;

/* 字体定义 */
--table-font: 'Microsoft YaHei', '微软雅黑', sans-serif;

/* 字号定义 */
--table-title-size: 28px;
--table-header-size: 20px;
--table-cell-size: 18px;

/* 间距定义 */
--table-title-spacing: 20px;
--table-header-height: 40px;
--table-cell-padding: 12px 16px;
--table-border-width: 1px;

/* 布局定义 */
--table-max-width: 800px;
--table-column-gap: 10px;
"""
    style_file.write_text(style_content, encoding='utf-8')

    # 更新 examples/example.py
    example_file = skill_path / 'examples' / 'example.py'
    example_content = """# 表格示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_table(pres, title, columns, rows, table_type='standard'):
    \"\"\"创建表格\"\"\"
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 56, 100)

    # 创建表格
    rows_count = len(rows) + 1
    cols_count = len(columns)

    table = slide.shapes.add_table(
        rows_count, cols_count, Inches(1), Inches(4), Inches(8), Inches(5)
    ).table

    # 设置列宽
    col_width = Inches(8) / cols_count
    for i in range(cols_count):
        table.columns[i].width = col_width

    # 表头
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(74, 111, 165)
        cell.text_frame.paragraphs[0].font.size = Pt(20)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 数据行
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(18)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 斑马纹效果
            if table_type == 'zebra' and i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 247, 250)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    columns = ['步骤', '操作', '时间', '负责人']
    rows = [
        ['1', '提交申请', '1-2天', '客户经理'],
        ['2', '审核材料', '2-3天', '风控专员'],
        ['3', '签订合同', '1天', '法务部门']
    ]
    slide = create_table(pres, '签约流程步骤对比', columns, rows, table_type='zebra')
    pres.save('output/table_example.pptx')
    print('表格示例已生成: output/table_example.pptx')
"""
    example_file.write_text(example_content, encoding='utf-8')

def enhance_flowchart_skill(skill_path):
    """增强流程图技能包"""
    print(f"  [+] 增强流程图技能包")

    # 更新 templates/template_example.txt
    template_file = skill_path / 'templates' / 'template_example.txt'
    template_content = """# 流程图模板

## 模板说明
标准流程图模板，展示业务流程和决策路径

## 使用方法
1. 替换 {{title}} 为流程图标题
2. 替换 {{nodes}} 为流程节点
3. 替换 {{edges}} 为节点连接关系
4. 选择流程图类型（标准/泳道/循环）

## 示例
{{title}}: 个人综合签约流程
{{nodes}}:
- 开始 -> 签约申请
- 签约申请 -> 材料审核
- 材料审核 -> 合同签订
- 合同签订 -> 完成
{{edges}}:
- 开始 -> 签约申请 [实线]
- 签约申请 -> 材料审核 [实线]
- 材料审核 -> 合同签订 [实线]
- 合同签订 -> 完成 [实线]
"""
    template_file.write_text(template_content, encoding='utf-8')

    # 更新 styles/style_example.css
    style_file = skill_path / 'styles' / 'style_example.css'
    style_content = """/* 流程图样式定义 */

/* 颜色定义 */
--flow-primary: #1F3864;
--flow-secondary: #2E75B6;
--flow-success: #4CAF50;
--flow-warning: #FF9800;
--flow-error: #F44336;
--flow-neutral: #9E9E9E;

/* 字体定义 */
--flow-font: 'Microsoft YaHei', '微软雅黑', sans-serif;

/* 字号定义 */
--flow-title-size: 28px;
--flow-node-size: 20px;
--flow-text-size: 18px;

/* 间距定义 */
--flow-title-spacing: 20px;
--flow-node-spacing: 40px;
--flow-node-padding: 15px 25px;

/* 布局定义 */
--flow-max-width: 900px;
--flow-node-radius: 8px;
--flow-line-width: 2px;
"""
    style_file.write_text(style_content, encoding='utf-8')

    # 更新 examples/example.py
    example_file = skill_path / 'examples' / 'example.py'
    example_content = """# 流程图示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_flowchart(pres, title, nodes, edges, flow_type='standard'):
    \"\"\"创建流程图\"\"\"
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 56, 100)

    # 创建流程节点
    start_x = Inches(2)
    start_y = Inches(2.5)
    node_width = Inches(1.5)
    node_height = Inches(0.6)

    # 绘制节点
    node_positions = []
    for i, node in enumerate(nodes):
        x = start_x + i * (node_width + Inches(0.5))
        y = start_y

        # 创建矩形形状
        shape = slide.shapes.add_shape(
            1,  # 矩形
            x, y, node_width, node_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(46, 117, 182)
        shape.line.color.rgb = RGBColor(31, 56, 100)

        # 设置圆角
        shape.shape_properties.preset = 13  # 圆角矩形

        # 添加文本
        text_frame = shape.text_frame
        text_frame.text = node
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        node_positions.append((x, y))

    # 绘制连接线
    for i in range(len(node_positions) - 1):
        x1, y1 = node_positions[i]
        x2, y2 = node_positions[i + 1]

        # 水平线
        line = slide.shapes.add_shape(
            1,  # 矩形
            x1 + node_width, y1 + node_height / 2 - 1,
            x2 - x1 - node_width, 2
        )
        line.line.color.rgb = RGBColor(31, 56, 100)
        line.line.width = Pt(2)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    nodes = ['开始', '签约申请', '材料审核', '合同签订', '完成']
    slide = create_flowchart(pres, '个人综合签约流程', nodes, [])
    pres.save('output/flowchart_example.pptx')
    print('流程图示例已生成: output/flowchart_example.pptx')
"""
    example_file.write_text(example_content, encoding='utf-8')

def enhance_scene_description_skill(skill_path):
    """增强场景描述技能包"""
    print(f"  [+] 增强场景描述技能包")

    # 更新 templates/template_example.txt
    template_file = skill_path / 'templates' / 'template_example.txt'
    template_content = """# 场景描述模板

## 模板说明
标准场景描述模板，展示业务场景和关键要素

## 使用方法
1. 替换 {{title}} 为场景标题
2. 替换 {{scenario}} 为场景描述
3. 替换 {{characters}} 为角色信息
4. 替换 {{environment}} 为环境信息

## 示例
{{title}}: 签约场景
{{scenario}}: 客户经理与客户在办公室进行签约
{{characters}}:
- 客户经理：负责引导和解释
- 客户：提交申请和确认信息
{{environment}}: 办公室环境，安静私密
"""
    template_file.write_text(template_content, encoding='utf-8')

    # 更新 styles/style_example.css
    style_file = skill_path / 'styles' / 'style_example.css'
    style_content = """/* 场景描述样式定义 */

/* 颜色定义 */
--scene-primary: #1F3864;
--scene-secondary: #2E75B6;
--scene-accent: #ED7D31;
--scene-text: #333333;
--scene-light: #F5F7FA;

/* 字体定义 */
--scene-title-font: 'Microsoft YaHei', '微软雅黑', sans-serif;
--scene-content-font: 'Microsoft YaHei', '微软雅黑', sans-serif;

/* 字号定义 */
--scene-title-size: 28px;
--scene-content-size: 20px;
--scene-label-size: 18px;

/* 间距定义 */
--scene-title-spacing: 20px;
--scene-content-spacing: 15px;
--scene-label-spacing: 10px;

/* 布局定义 */
--scene-max-width: 800px;
--scene-card-padding: 20px;
"""
    style_file.write_text(style_content, encoding='utf-8')

    # 更新 examples/example.py
    example_file = skill_path / 'examples' / 'example.py'
    example_content = """# 场景描述示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_scene_description(pres, title, scenario, characters, environment):
    \"\"\"创建场景描述页\"\"\"
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 56, 100)

    # 场景描述
    y_position = Inches(1.5)
    x_position = Inches(2)

    # 场景描述框
    desc_box = slide.shapes.add_textbox(
        x_position, y_position, Inches(6), Inches(1.5)
    )
    desc_frame = desc_box.text_frame
    desc_frame.text = scenario

    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(20)
    desc_para.font.color.rgb = RGBColor(51, 51, 51)
    desc_para.alignment = PP_ALIGN.CENTER

    # 角色信息
    y_position += Inches(2)
    label_box = slide.shapes.add_textbox(
        x_position, y_position, Inches(2), Inches(0.5)
    )
    label_frame = label_box.text_frame
    label_frame.text = '角色：'
    label_frame.paragraphs[0].font.size = Pt(18)
    label_frame.paragraphs[0].font.bold = True
    label_frame.paragraphs[0].font.color.rgb = RGBColor(31, 56, 100)

    chars_box = slide.shapes.add_textbox(
        x_position + Inches(2), y_position, Inches(4), Inches(0.5)
    )
    chars_frame = chars_box.text_frame
    chars_frame.text = characters

    chars_para = chars_frame.paragraphs[0]
    chars_para.font.size = Pt(18)
    chars_para.font.color.rgb = RGBColor(51, 51, 51)

    # 环境信息
    y_position += Inches(0.8)
    label_box2 = slide.shapes.add_textbox(
        x_position, y_position, Inches(2), Inches(0.5)
    )
    label_frame2 = label_box2.text_frame
    label_frame2.text = '环境：'
    label_frame2.paragraphs[0].font.size = Pt(18)
    label_frame2.paragraphs[0].font.bold = True
    label_frame2.paragraphs[0].font.color.rgb = RGBColor(31, 56, 100)

    env_box = slide.shapes.add_textbox(
        x_position + Inches(2), y_position, Inches(4), Inches(0.5)
    )
    env_frame = env_box.text_frame
    env_frame.text = environment

    env_para = env_frame.paragraphs[0]
    env_para.font.size = Pt(18)
    env_para.font.color.rgb = RGBColor(51, 51, 51)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    slide = create_scene_description(
        pres,
        title='签约场景',
        scenario='客户经理与客户在办公室进行签约',
        characters='客户经理：负责引导和解释\n客户：提交申请和确认信息',
        environment='办公室环境，安静私密'
    )
    pres.save('output/scene_description_example.pptx')
    print('场景描述示例已生成: output/scene_description_example.pptx')
"""
    example_file.write_text(example_content, encoding='utf-8')

def main():
    """主函数"""
    print("=" * 60)
    print("技能包内容增强工具")
    print("=" * 60)

    total_enhanced = 0

    for skill_name in SKILL_PACKAGES:
        skill_path = Path('skills') / skill_name
        if not skill_path.exists():
            print(f"\n技能包不存在: {skill_name}")
            continue

        print(f"\n处理技能包: {skill_name}")

        # 根据技能包类型进行增强
        if skill_name == 'cover-skill':
            enhance_cover_skill(skill_path)
            total_enhanced += 1
        elif skill_name == 'toc-skill':
            enhance_toc_skill(skill_path)
            total_enhanced += 1
        elif skill_name == 'table-skill':
            enhance_table_skill(skill_path)
            total_enhanced += 1
        elif skill_name == 'flowchart-skill':
            enhance_flowchart_skill(skill_path)
            total_enhanced += 1
        elif skill_name == 'scene-description-skill':
            enhance_scene_description_skill(skill_path)
            total_enhanced += 1
        else:
            # 其他技能包暂时使用默认模板
            print(f"  [i] 使用默认模板（暂不增强）")

    print("\n" + "=" * 60)
    print("内容增强完成")
    print("=" * 60)
    print(f"已增强技能包数量: {total_enhanced}")

if __name__ == '__main__':
    main()
