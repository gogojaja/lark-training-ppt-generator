#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""generate_page.py — 基于框架的端到端页面生成器

将 style-brief 风格配方 + vertical_rhythm 垂直节奏 + 角色字号驱动，
渲染为真实 PPTX。支持多种页面类型，可用于整本文档生成与闭环验证。

== 依赖 ==
    pip install python-pptx

== 用法 ==
    # 单页（默认 content 类型，使用示例内容）
    py -3 generate_page.py --out 生成产物/demo_业务要点.pptx

    # 多页 deck（从 JSON 读取页面列表）
    py -3 generate_page.py --deck my_deck.json --out 生成产物/demo_deck.pptx

    # 指定风格配方
    py -3 generate_page.py --recipe ../skills/style-brief-skill/recipes/red-alert.json

== content / deck JSON 结构 ==
单页: { "type": "content|cover|toc|steps|warning", ... }
deck: { "pages": [ <单页>, ... ] }
"""
import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.util import Cm, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(SCRIPT_DIR, "..", "skills", "ppt-framework"))
from vertical_rhythm import VerticalRhythm  # noqa: E402

SLIDE_W_CM = 33.867
SLIDE_H_CM = 19.05

ROLE_MAP = {
    "sectionLabel": {"size": 12, "bold": False, "color": "secondary", "lh": 1.3},
    "claim":        {"size": 32, "bold": True,  "color": "primary",   "lh": 1.2},
    "body":         {"size": 14, "bold": False, "color": "text",      "lh": 1.5},
    "annotation":   {"size": 12, "bold": False, "color": "text",      "lh": 1.4},
    "source":       {"size": 10, "bold": False, "color": "muted",     "lh": 1.3},
}


def hex2rgb(hexstr):
    hexstr = hexstr.lstrip("#")
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


def blend_hex(c1, c2, w):
    """线性混合两个 HEX 颜色，w 为 c2 的权重（0..1）。用于生成柔和浅色。"""
    def _ch(c):
        c = c.lstrip("#")
        return [int(c[i:i + 2], 16) for i in (0, 2, 4)]
    a, b = _ch(c1), _ch(c2)
    out = [round(a[i] * (1 - w) + b[i] * w) for i in range(3)]
    return "%02X%02X%02X" % tuple(min(255, max(0, v)) for v in out)


class PageBuilder:
    def __init__(self, recipe):
        self.recipe = recipe
        pal = recipe.get("palette", {})
        self.colors = {
            "primary": pal.get("primary", {}).get("value", "1F3864"),
            "secondary": pal.get("secondary", {}).get("value", "2E75B6"),
            "accent": pal.get("accent", {}).get("value", "ED7D31"),
            "text": pal.get("text", {}).get("value", "333333"),
            "muted": "999999",
            "background": pal.get("background", {}).get("value", "FFFFFF"),
            "card": "F8FAFC",
            "cardBorder": "E2E8F0",
            "light": "E8EEF7",
        }
        self.vr = VerticalRhythm()
        self.layout = self.vr.get_layout("standard")

    # ---------- 基础工具 ----------
    def _band_cm(self, name):
        b = self.layout.get_band(name)
        return b.start_cm, b.height_cm

    def _set_text(self, tf, text, role, align=PP_ALIGN.LEFT, color_override=None, size_override=None):
        cfg = ROLE_MAP[role]
        color = self.colors.get(color_override or cfg["color"], "333333")
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        f = run.font
        f.size = Pt(size_override or cfg["size"])
        f.bold = cfg["bold"]
        f.name = "Microsoft YaHei"
        f.color.rgb = hex2rgb(color)
        p.line_spacing = cfg["lh"]

    def _new_slide(self, prs, bg=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = hex2rgb(bg or self.colors["background"])
        return slide

    def _footer(self, slide, source, page_no):
        y4, _ = self._band_cm("footer")
        ftb = slide.shapes.add_textbox(Cm(2.5), Cm(y4 + 0.1), Cm(SLIDE_W_CM - 7.5), Cm(1.2))
        self._set_text(ftb.text_frame, source, "source")
        ptb = slide.shapes.add_textbox(Cm(SLIDE_W_CM - 5), Cm(y4 + 0.1), Cm(2.5), Cm(1.2))
        self._set_text(ptb.text_frame, page_no, "source", align=PP_ALIGN.RIGHT)

    def _title_block(self, slide, context, claim, dark=False):
        """上下文 + 主张 + 短分隔线（标准内容页标题区）"""
        y0, _ = self._band_cm("context")
        tb = slide.shapes.add_textbox(Cm(2.5), Cm(y0 + 0.2), Cm(28), Cm(1.3))
        self._set_text(tb.text_frame, context, "sectionLabel",
                       color_override=("muted" if dark else "secondary"))
        y1, h1 = self._band_cm("claim")
        tb = slide.shapes.add_textbox(Cm(2.5), Cm(y1), Cm(28), Cm(h1))
        tf = tb.text_frame
        tf.word_wrap = True
        self._set_text(tf, claim, "claim",
                       color_override=("background" if dark else "primary"))
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Cm(2.5), Cm(y1 + h1 - 0.15), Cm(6), Cm(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = hex2rgb(self.colors["accent"])
        line.line.fill.background()

    # ---------- 各页面类型 ----------
    def build_content(self, slide, content):
        self._title_block(slide, content.get("context", ""), content.get("claim", ""))
        y2, h2 = self._band_cm("evidence")
        cards = content.get("cards", [])
        n = len(cards)
        cols = 1 if n == 1 else (2 if n <= 2 else 3)
        margin_x = 2.5
        gap = 0.6
        total_w = SLIDE_W_CM - margin_x * 2
        card_w = (total_w - gap * (cols - 1)) / cols
        card_h = min((h2 - gap) / 1.0, 5.2)
        start_y = y2 + 0.6
        for i, card in enumerate(cards):
            x = margin_x + (i % cols) * (card_w + gap)
            y = start_y + (i // cols) * (card_h + gap)
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Cm(x), Cm(y), Cm(card_w), Cm(card_h))
            panel.fill.solid()
            panel.fill.fore_color.rgb = hex2rgb(self.colors["card"])
            panel.line.color.rgb = hex2rgb(self.colors["cardBorder"])
            panel.line.width = Cm(0.03)
            panel.shadow.inherit = False
            ctf = panel.text_frame
            ctf.word_wrap = True
            ctf.margin_left = Cm(0.4)
            ctf.margin_right = Cm(0.4)
            ctf.margin_top = Cm(0.3)
            ctf.vertical_anchor = MSO_ANCHOR.TOP
            pt = ctf.paragraphs[0]
            run = pt.add_run()
            run.text = card.get("title", "")
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = hex2rgb(self.colors["primary"])
            pb = ctf.add_paragraph()
            run = pb.add_run()
            run.text = card.get("body", "")
            run.font.size = Pt(11)
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = hex2rgb(self.colors["text"])
            pb.line_spacing = 1.4
            pb.space_before = Cm(0.15)
        # 含义收尾条
        y3, h3 = self._band_cm("meaning")
        strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Cm(2.5), Cm(y3 + 0.15),
                                       Cm(SLIDE_W_CM - 5), Cm(h3 - 0.3))
        strip.fill.solid()
        strip.fill.fore_color.rgb = hex2rgb(self.colors["primary"])
        strip.line.fill.background()
        strip.shadow.inherit = False
        stf = strip.text_frame
        stf.word_wrap = True
        stf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = stf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "要点  |  " + content.get("meaning", "")
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.name = "Microsoft YaHei"
        run.font.color.rgb = hex2rgb("FFFFFF")
        self._footer(slide, content.get("source", "来源：内部培训材料"),
                     content.get("page", "01 / 01"))

    def build_cover(self, slide, content):
        primary = self.colors["primary"]
        # 柔和浅色背景：主色向白色混合约 88%（柔和粉彩），避免深黑/刺眼
        soft_bg = blend_hex(primary, "FFFFFF", 0.88)
        # 柔和强调带：主色向白色混合约 72%，用于横幅/色块
        soft_band = blend_hex(primary, "FFFFFF", 0.72)
        # 标题色：略加深的主色，保证在浅色背景上的对比度与柔和感
        title_color = blend_hex(primary, "000000", 0.10)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = hex2rgb(soft_bg)
        # 顶部柔和色带（浅色横幅，替代刺眼的整页深色）
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Cm(0), Cm(0), Cm(SLIDE_W_CM), Cm(3.2))
        band.fill.solid()
        band.fill.fore_color.rgb = hex2rgb(soft_band)
        band.line.fill.background()
        # 主标题（居中）
        tb = slide.shapes.add_textbox(Cm(4), Cm(7.2), Cm(26), Cm(4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = content.get("claim", "")
        run.font.size = Pt(38)
        run.font.bold = True
        run.font.name = "Microsoft YaHei"
        run.font.color.rgb = hex2rgb(title_color)
        # 副标题
        st = slide.shapes.add_textbox(Cm(4), Cm(12.2), Cm(26), Cm(2))
        tf2 = st.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = content.get("context", "")
        run.font.size = Pt(16)
        run.font.name = "Microsoft YaHei"
        run.font.color.rgb = hex2rgb(blend_hex(primary, "FFFFFF", 0.55))
        # 居中短分隔线（柔和强调色，细线不刺眼）
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Cm(14.4), Cm(11.5), Cm(5), Cm(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = hex2rgb(blend_hex(self.colors["accent"], "FFFFFF", 0.30))
        line.line.fill.background()

    def build_toc(self, slide, content):
        items = content.get("items", [])
        n = len(items)
        raw = (content.get("claim") or "").strip()
        toc_claim = raw if raw and raw != "目录" else f"目录（共 {n} 项）"
        self._title_block(slide, content.get("context", "培训结构总览"), toc_claim)
        y2, h2 = self._band_cm("evidence")
        items = content.get("items", [])
        n = len(items)
        margin_x = 2.5
        row_h = min((h2 - 0.6) / max(n, 1), 2.6)
        for i, item in enumerate(items):
            y = y2 + 0.6 + i * row_h
            # 编号徽章
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                           Cm(margin_x), Cm(y + 0.2),
                                           Cm(1.1), Cm(1.1))
            badge.fill.solid()
            badge.fill.fore_color.rgb = hex2rgb(self.colors["primary"])
            badge.line.fill.background()
            badge.shadow.inherit = False
            btf = badge.text_frame
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = btf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(i + 1)
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = hex2rgb("FFFFFF")
            # 条目文本
            ttb = slide.shapes.add_textbox(Cm(margin_x + 1.6), Cm(y + 0.15),
                                           Cm(SLIDE_W_CM - margin_x - 4), Cm(row_h - 0.3))
            tf = ttb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = hex2rgb(self.colors["text"])
        self._footer(slide, content.get("source", "来源：内部培训材料"),
                     content.get("page", "01 / 01"))

    def build_steps(self, slide, content):
        self._title_block(slide, content.get("context", ""), content.get("claim", ""))
        y2, h2 = self._band_cm("evidence")
        steps = content.get("steps", [])
        n = len(steps)
        margin_x = 2.5
        row_h = min((h2 - 0.6) / max(n, 1), 4.0)
        for i, step in enumerate(steps):
            y = y2 + 0.6 + i * row_h
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                            Cm(margin_x), Cm(y + 0.2),
                                            Cm(1.3), Cm(1.3))
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex2rgb(self.colors["accent"])
            circle.line.fill.background()
            circle.shadow.inherit = False
            ctf = circle.text_frame
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = ctf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(i + 1)
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = hex2rgb("FFFFFF")
            # 标题 + 说明
            ttb = slide.shapes.add_textbox(Cm(margin_x + 1.8), Cm(y + 0.1),
                                           Cm(SLIDE_W_CM - margin_x - 4.3), Cm(row_h - 0.2))
            tf = ttb.text_frame
            tf.word_wrap = True
            pt = tf.paragraphs[0]
            run = pt.add_run()
            run.text = step.get("title", "")
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = hex2rgb(self.colors["primary"])
            if step.get("body"):
                pb = tf.add_paragraph()
                run = pb.add_run()
                run.text = step.get("body", "")
                run.font.size = Pt(11)
                run.font.name = "Microsoft YaHei"
                run.font.color.rgb = hex2rgb(self.colors["text"])
                pb.line_spacing = 1.4
                pb.space_before = Cm(0.1)
        self._footer(slide, content.get("source", "来源：内部培训材料"),
                     content.get("page", "01 / 01"))

    def build_warning(self, slide, content):
        self._title_block(slide, content.get("context", ""), content.get("claim", ""))
        y2, h2 = self._band_cm("evidence")
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Cm(2.5), Cm(y2 + 0.6),
                                       Cm(SLIDE_W_CM - 5), Cm(h2 - 1.0))
        panel.fill.solid()
        panel.fill.fore_color.rgb = hex2rgb("FDF1E7")
        panel.line.color.rgb = hex2rgb(self.colors["accent"])
        panel.line.width = Cm(0.06)
        panel.shadow.inherit = False
        ctf = panel.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Cm(0.6)
        ctf.margin_right = Cm(0.6)
        ctf.margin_top = Cm(0.4)
        ctf.vertical_anchor = MSO_ANCHOR.TOP
        pt = ctf.paragraphs[0]
        run = pt.add_run()
        run.text = "⚠ " + content.get("warning_title", "注意事项")
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.name = "Microsoft YaHei"
        run.font.color.rgb = hex2rgb(self.colors["accent"])
        for line_text in content.get("points", []):
            pb = ctf.add_paragraph()
            run = pb.add_run()
            run.text = "• " + line_text
            run.font.size = Pt(13)
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = hex2rgb(self.colors["text"])
            pb.line_spacing = 1.5
            pb.space_before = Cm(0.15)
        self._footer(slide, content.get("source", "来源：内部培训材料"),
                     content.get("page", "01 / 01"))

    # ---------- 调度 ----------
    BUILDERS = {
        "content": build_content,
        "cover": build_cover,
        "toc": build_toc,
        "steps": build_steps,
        "warning": build_warning,
    }

    def build(self, content, out_path):
        prs = Presentation()
        prs.slide_width = Emu(int(SLIDE_W_CM * 360000))
        prs.slide_height = Emu(int(SLIDE_H_CM * 360000))
        ptype = content.get("type", "content")
        slide = self._new_slide(prs)
        self.BUILDERS.get(ptype, self.build_content)(self, slide, content)
        prs.save(out_path)
        return out_path

    def build_deck(self, pages, out_path):
        prs = Presentation()
        prs.slide_width = Emu(int(SLIDE_W_CM * 360000))
        prs.slide_height = Emu(int(SLIDE_H_CM * 360000))
        total = len(pages)
        for idx, content in enumerate(pages, 1):
            content = dict(content)
            if "page" not in content:
                content["page"] = f"{idx:02d} / {total:02d}"
            slide = self._new_slide(prs)
            ptype = content.get("type", "content")
            self.BUILDERS.get(ptype, self.build_content)(self, slide, content)
        prs.save(out_path)
        return out_path


def load_recipe(path):
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def default_deck():
    return {
        "pages": [
            {"type": "cover", "claim": "个人综合签约 · 操作培训（2026版）",
             "context": "柜员 / 智能柜员机 / Pad 三端统一签约能力宣讲"},
            {"type": "toc", "claim": "目录",
             "items": ["业务概述", "签约规则", "操作步骤", "注意事项"]},
            {"type": "content", "context": "业务要点 · 个人综合签约",
             "claim": "签约费每月2元，扣费失败系统25日自动补扣",
             "cards": [
                 {"title": "计费规则", "body": "系统每月15日自动扣款，每户每月2元/月签约费，统一托收。"},
                 {"title": "失败重试", "body": "首次扣费失败，于25日再次自动扣款，无需柜员手工干预。"},
                 {"title": "授权要求", "body": "业务代码036101，需业务主管授权方可办理签约与解约。"},
             ],
             "meaning": "签约费由系统统一托收，柜员无需手工收费，降低操作风险。",
             "source": "来源：个人综合签约操作手册 v1.2"},
            {"type": "steps", "context": "操作步骤 · 智能柜员机",
             "claim": "四步完成客户综合签约",
             "steps": [
                 {"title": "身份核验", "body": "刷脸 + 身份证读取，确认客户本人。"},
                 {"title": "选择产品", "body": "在签约列表勾选「个人综合签约」。"},
                 {"title": "主管授权", "body": "业务主管 UKey 授权通过。"},
                 {"title": "回单打印", "body": "系统生成电子回单，客户签字确认。"},
             ],
             "source": "来源：个人综合签约操作手册 v1.2"},
            {"type": "warning", "context": "注意事项 · 签约解约",
             "claim": "三类情形必须先解约前置签约",
             "warning_title": "红线提醒",
             "points": [
                 "签约账户销户前，必须先解除综合签约及其前置绑定关系。",
                 "扣费失败连续两月，需主动联系客户更新缴费账户。",
                 "授权权限仅限业务主管，严禁柜员越权办理。",
             ],
             "source": "来源：个人综合签约操作手册 v1.2"},
        ]
    }


def main():
    ap = argparse.ArgumentParser(description="框架驱动的 PPT 页面生成器")
    ap.add_argument("--recipe", default=os.path.join(
        SCRIPT_DIR, "..", "skills", "style-brief-skill", "recipes", "professional-blue.json"))
    ap.add_argument("--content", help="单页 content JSON 路径")
    ap.add_argument("--deck", help="多页 deck JSON 路径（含 pages 列表）")
    ap.add_argument("--out", default=os.path.join(SCRIPT_DIR, "..", "生成产物", "demo_deck.pptx"))
    args = ap.parse_args()

    recipe = load_recipe(args.recipe)
    builder = PageBuilder(recipe)

    if args.deck:
        with open(args.deck, encoding="utf-8") as f:
            data = json.load(f)
        pages = data.get("pages", [data])
        out = builder.build_deck(pages, args.out)
    elif args.content:
        with open(args.content, encoding="utf-8") as f:
            content = json.load(f)
        out = builder.build(content, args.out)
    else:
        out = builder.build_deck(default_deck()["pages"], args.out)
    print("已生成:", os.path.abspath(out))


if __name__ == "__main__":
    main()
