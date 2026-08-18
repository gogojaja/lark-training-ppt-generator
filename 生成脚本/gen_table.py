#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
表格页生成脚本
支持多种模板和配色方案
"""

import sys
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 颜色方案
COLOR_SCHEMES = {
    'blue': {
        'header_bg': RGBColor(31, 56, 100),      # 1F3864
        'header_text': RGBColor(255, 255, 255),  # FFFFFF
        'row_bg': RGBColor(255, 255, 255),       # FFFFFF
        'zebra_bg': RGBColor(248, 250, 252),     # F8FAFC
        'border': RGBColor(224, 224, 224),       # E0E0E0
        'accent': RGBColor(46, 117, 182),        # 2E75B6
        'text': RGBColor(51, 51, 51)             # 333333
    },
    'gray': {
        'header_bg': RGBColor(64, 64, 64),       # 404040
        'header_text': RGBColor(255, 255, 255),  # FFFFFF
        'row_bg': RGBColor(255, 255, 255),       # FFFFFF
        'zebra_bg': RGBColor(245, 245, 245),     # F5F5F5
        'border': RGBColor(224, 224, 224),       # E0E0E0
        'accent': RGBColor(128, 128, 128),       # 808080
        'text': RGBColor(51, 51, 51)             # 333333
    }
}

TEMPLATE_TYPES = ['standard', 'zebra', 'icon']

def create_table_slide(pres, data, template='standard', preset='blue', title=None):
    """
    创建表格页

    Args:
        pres: PPT presentation object
        data: 表格数据字典，包含 title, headers, rows
        template: 模板类型
        preset: 配色方案
        title: 自定义标题
    """
    colors = COLOR_SCHEMES.get(preset, COLOR_SCHEMES['blue'])
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 设置背景
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colors['row_bg']

    # 标题
    title_text = title or data.get('title', '表格')
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = 'Microsoft YaHei'
    run.font.color.rgb = colors['header_bg']

    # 计算表格尺寸
    headers = data.get('headers', [])
    rows = data.get('rows', [])
    num_cols = len(headers)
    num_rows = len(rows) + 1  # 包括表头

    # 表格宽度：留出边距
    table_width = Inches(8.5)
    margin = Inches(0.5)
    table_x = margin
    table_y = Inches(1.5)

    # 创建表格
    table = slide.shapes.add_table(num_rows, num_cols, table_x, table_y, table_width, Inches(4.5)).table

    # 设置列宽（平均分配）
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # 填充表头
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors['header_bg']
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 表头文字
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.name = 'Microsoft YaHei'
                run.font.color.rgb = colors['header_text']

    # 填充数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.rows[row_idx + 1].cells[col_idx]
            cell.text = str(cell_data)

            # 设置文字格式
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.name = 'Microsoft YaHei'
                    run.font.color.rgb = colors['text']

            # 斑马纹效果
            if template == 'zebra' and row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors['zebra_bg']

    # 设置边框
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    pass

    # 添加页脚
    add_footer(slide, '表格页')

    return slide

def add_footer(slide, text):
    """添加页脚"""
    y = Inches(6.0)
    footer = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.font.size = Pt(10)
    run.font.name = 'Microsoft YaHei'
    run.font.color.rgb = RGBColor(153, 153, 153)

def main():
    """主函数"""
    if len(sys.argv) < 2:
        print('Usage: python gen_table.py <json_file> [template] [preset] [output_file]')
        print('Templates: standard, zebra, icon')
        print('Presets: blue, gray')
        sys.exit(1)

    json_file = sys.argv[1]
    template = sys.argv[2] if len(sys.argv) > 2 else 'standard'
    preset = sys.argv[3] if len(sys.argv) > 3 else 'blue'
    output_file = sys.argv[4] if len(sys.argv) > 4 else '表格.pptx'

    # 检查模板和预设是否有效
    if template not in TEMPLATE_TYPES:
        print(f'Error: Invalid template. Must be one of {TEMPLATE_TYPES}')
        sys.exit(1)

    if preset not in COLOR_SCHEMES:
        print(f'Error: Invalid preset. Must be one of {list(COLOR_SCHEMES.keys())}')
        sys.exit(1)

    # 读取 JSON 数据
    with open(json_file, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # 创建 PPT
    pres = Presentation()
    pres.slide_width = int(Inches(9.5))
    pres.slide_height = int(Inches(6.5))

    # 创建表格页
    create_table_slide(pres, data, template, preset)

    # 保存
    output_path = output_file
    pres.save(output_path)
    print(f'Saved to: {output_path}')

if __name__ == '__main__':
    main()
