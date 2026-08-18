# 场景描述示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_scene_description(pres, title, scenario, characters, environment):
    """创建场景描述页"""
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 56, 100)

    # 场景描述
    y_position = Inches(1.5)
    x_position = Inches(2)

    # 场景描述框
    desc_box = slide.shapes.add_textbox(
        x_position, y_position, Inches(6), Inches(1.5)
    )
    desc_frame = desc_box.text_frame
    desc_frame.text = scenario

    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(20)
    desc_para.font.color.rgb = RGBColor(51, 51, 51)
    desc_para.alignment = PP_ALIGN.CENTER

    # 角色信息
    y_position += Inches(2)
    label_box = slide.shapes.add_textbox(
        x_position, y_position, Inches(2), Inches(0.5)
    )
    label_frame = label_box.text_frame
    label_frame.text = '角色：'
    label_frame.paragraphs[0].font.size = Pt(18)
    label_frame.paragraphs[0].font.bold = True
    label_frame.paragraphs[0].font.color.rgb = RGBColor(31, 56, 100)

    chars_box = slide.shapes.add_textbox(
        x_position + Inches(2), y_position, Inches(4), Inches(0.5)
    )
    chars_frame = chars_box.text_frame
    chars_frame.text = characters

    chars_para = chars_frame.paragraphs[0]
    chars_para.font.size = Pt(18)
    chars_para.font.color.rgb = RGBColor(51, 51, 51)

    # 环境信息
    y_position += Inches(0.8)
    label_box2 = slide.shapes.add_textbox(
        x_position, y_position, Inches(2), Inches(0.5)
    )
    label_frame2 = label_box2.text_frame
    label_frame2.text = '环境：'
    label_frame2.paragraphs[0].font.size = Pt(18)
    label_frame2.paragraphs[0].font.bold = True
    label_frame2.paragraphs[0].font.color.rgb = RGBColor(31, 56, 100)

    env_box = slide.shapes.add_textbox(
        x_position + Inches(2), y_position, Inches(4), Inches(0.5)
    )
    env_frame = env_box.text_frame
    env_frame.text = environment

    env_para = env_frame.paragraphs[0]
    env_para.font.size = Pt(18)
    env_para.font.color.rgb = RGBColor(51, 51, 51)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    slide = create_scene_description(
        pres,
        title='签约场景',
        scenario='客户经理与客户在办公室进行签约',
        characters='客户经理：负责引导和解释
客户：提交申请和确认信息',
        environment='办公室环境，安静私密'
    )
    pres.save('output/scene_description_example.pptx')
    print('场景描述示例已生成: output/scene_description_example.pptx')
