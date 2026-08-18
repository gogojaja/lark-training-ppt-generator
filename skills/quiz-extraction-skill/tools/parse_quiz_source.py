#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""文档解析工具（parse_quiz_source.py）

考试题抽取技能 Step 1：把 Word/PPT/TXT/Markdown 源文档解析为统一中间文本结构，
供 AI 语义抽取使用。解析层只做机械提取，不做语义判断。

用法（Windows / macOS 双平台）:
    python tools/parse_quiz_source.py 输入文档/试卷.docx -o 输出/中间结构.json
    python tools/parse_quiz_source.py 输入文档/课件.pptx -o 输出/中间结构.json
    python tools/parse_quiz_source.py 输入文档/试题.txt -o 输出/中间结构.json

输出:
    {
      "paper_title": "...",
      "format": "docx",
      "pages": [...],
      "blocks": [
        {"type": "paragraph|table", "text": "...", "page": 1, "order": 0},
        ...
      ]
    }
"""
import argparse
import json
import os
import sys


def parse_docx(path: str) -> dict:
    import docx
    doc = docx.Document(path)
    blocks = []
    order = 0
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            blocks.append({"type": "paragraph", "text": text, "page": None, "order": order})
            order += 1
    for tbl in doc.tables:
        rows = []
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            blocks.append({"type": "table", "text": "\n".join(rows), "page": None, "order": order})
            order += 1
    return {"paper_title": os.path.basename(path), "format": "docx",
            "pages": len(doc.sections), "blocks": blocks}


def parse_pptx(path: str) -> dict:
    from pptx import Presentation
    prs = Presentation(path)
    blocks = []
    order = 0
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    row_text = " | ".join(cells)
                    if row_text:
                        texts.append(row_text)
        if texts:
            blocks.append({"type": "paragraph", "text": "\n".join(texts), "page": idx, "order": order})
            order += 1
    return {"paper_title": os.path.basename(path), "format": "pptx",
            "pages": len(prs.slides), "blocks": blocks}


def parse_txt(path: str) -> dict:
    fmt = "md" if path.lower().endswith((".md", ".markdown")) else "txt"
    with open(path, encoding="utf-8") as fh:
        lines = fh.read().splitlines()
    blocks = []
    for order, line in enumerate(lines):
        if line.strip():
            blocks.append({"type": "paragraph", "text": line.strip(), "page": None, "order": order})
    return {"paper_title": os.path.basename(path), "format": fmt,
            "pages": 1, "blocks": blocks}


def main():
    ap = argparse.ArgumentParser(description="考试题抽取：源文档解析为中间文本结构")
    ap.add_argument("input", help="源文档路径（docx/pptx/txt/md）")
    ap.add_argument("-o", "--out", default="中间结构.json", help="输出 JSON 路径")
    args = ap.parse_args()

    path = args.input
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        data = parse_docx(path)
    elif ext == ".pptx":
        data = parse_pptx(path)
    elif ext in (".txt", ".md", ".markdown"):
        data = parse_txt(path)
    else:
        print(f"[错误] 不支持的格式: {ext}（仅支持 .docx/.pptx/.txt/.md）")
        sys.exit(1)

    data["blocks_count"] = len(data["blocks"])
    os.makedirs(os.path.dirname(os.path.abspath(args.out)), exist_ok=True)
    with open(args.out, "w", encoding="utf-8") as fh:
        json.dump(data, fh, ensure_ascii=False, indent=2)
    print(f"[完成] 解析 {path} -> {args.out}（{len(data['blocks'])} 个文本块）")


if __name__ == "__main__":
    main()