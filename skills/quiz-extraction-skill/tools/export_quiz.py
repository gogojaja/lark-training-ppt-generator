#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""多格式输出工具（export_quiz.py）

考试题抽取技能 Step 4：把规范化 JSON 输出为 CSV（UTF-8 BOM，Excel 直接打开）、
Markdown 试卷、PPTX 问卷。

用法:
    python tools/export_quiz.py csv  输出/规范化.json -o 输出/题库.csv
    python tools/export_quiz.py md   输出/规范化.json -o 输出/试卷.md
    python tools/export_quiz.py pptx 输出/规范化.json -o 输出/问卷.pptx
    # 缺省一次性输出全部
    python tools/export_quiz.py all  输出/规范化.json -o 输出目录/

CSV 列:
    group_id,section,shared_stem,qid,type,stem,option_a,option_b,option_c,
    option_d,option_e,option_f,answer,analysis,difficulty,score,source_page
"""
import argparse
import json
import os
import re
import sys

CSV_HEADER = ["group_id", "section", "shared_stem", "qid", "type", "stem",
              "option_a", "option_b", "option_c", "option_d", "option_e", "option_f",
              "answer", "analysis", "difficulty", "score", "source_page"]

TYPE_LABEL = {
    "single_choice": "单选题",
    "multiple_choice": "多选题",
    "true_false": "判断题",
    "fill_blank": "填空题",
    "short_answer": "简答题",
}


def _opt(q, i):
    opts = q.get("options") or []
    return opts[i] if i < len(opts) else ""


def _row(g, q):
    return {
        "group_id": g.get("group_id", ""),
        "section": g.get("section", ""),
        "shared_stem": g.get("shared_stem", "") or "",
        "qid": q.get("qid", ""),
        "type": q.get("type", ""),
        "stem": q.get("stem", ""),
        "option_a": _opt(q, 0), "option_b": _opt(q, 1), "option_c": _opt(q, 2),
        "option_d": _opt(q, 3), "option_e": _opt(q, 4), "option_f": _opt(q, 5),
        "answer": "/".join(q.get("answer") or []),
        "analysis": q.get("analysis", ""),
        "difficulty": q.get("difficulty", ""),
        "score": q.get("score", ""),
        "source_page": q.get("source_page", ""),
    }


def _iter_q(data):
    for g in data.get("groups", []):
        for q in g.get("questions", []):
            yield g, q


def export_csv(data, out):
    import csv
    os.makedirs(os.path.dirname(os.path.abspath(out)), exist_ok=True)
    with open(out, "w", encoding="utf-8-sig", newline="") as fh:
        w = csv.DictWriter(fh, fieldnames=CSV_HEADER)
        w.writeheader()
        for g, q in _iter_q(data):
            w.writerow(_row(g, q))
    print(f"[完成] CSV -> {out}")


def export_md(data, out):
    os.makedirs(os.path.dirname(os.path.abspath(out)), exist_ok=True)
    lines = [f"# {data.get('paper_title', '试卷')}", ""]
    for g in data.get("groups", []):
        if g.get("section"):
            lines.append(f"## {g['section']}")
            lines.append("")
        if g.get("shared_stem"):
            lines.append("**材料**：")
            lines.append("")
            lines.append(g["shared_stem"])
            lines.append("")
        for q in g.get("questions", []):
            label = TYPE_LABEL.get(q.get("type"), q.get("type", ""))
            lines.append(f"**{q.get('qid')}（{label}）** {q.get('stem')}")
            lines.append("")
            for opt in q.get("options") or []:
                lines.append(f"- {opt}")
                lines.append("")
            lines.append(f"> 答案：{('/'.join(q.get('answer') or [''])) or '待定'}")
            if q.get("analysis"):
                lines.append(f"> 解析：{q['analysis']}")
            lines.append("")
    md = "\n".join(lines)
    with open(out, "w", encoding="utf-8") as fh:
        fh.write(md)
    print(f"[完成] Markdown -> {out}")


def export_pptx(data, out):
    from pptx import Presentation
    from pptx.util import Pt, Inches

    prs = Presentation()
    for g in data.get("groups", []):
        if g.get("section"):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = g["section"]
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = ""
            if g.get("shared_stem"):
                p = tf.add_paragraph()
                p.text = f"材料：{g['shared_stem']}"
                p.font.size = Pt(14)
            for q in g.get("questions", []):
                label = TYPE_LABEL.get(q.get("type"), q.get("type", ""))
                p = tf.add_paragraph()
                p.text = f"{q.get('qid')}（{label}）{q.get('stem')}"
                p.font.size = Pt(16)
                for opt in q.get("options") or []:
                    po = tf.add_paragraph()
                    po.text = opt
                    po.font.size = Pt(14)
    os.makedirs(os.path.dirname(os.path.abspath(out)), exist_ok=True)
    prs.save(out)
    print(f"[完成] PPTX -> {out}")


def main():
    ap = argparse.ArgumentParser(description="考试题抽取：多格式输出")
    ap.add_argument("fmt", choices=["csv", "md", "pptx", "all"])
    ap.add_argument("input", help="规范化 JSON 路径")
    ap.add_argument("-o", "--out", help="输出路径或目录")
    args = ap.parse_args()

    with open(args.input, encoding="utf-8") as fh:
        data = json.load(fh)

    base, ext = os.path.splitext(args.input)
    out = args.out or base

    if args.fmt == "csv":
        export_csv(data, out if out.endswith(".csv") else out + ".csv")
    elif args.fmt == "md":
        export_md(data, out if out.endswith(".md") else out + ".md")
    elif args.fmt == "pptx":
        export_pptx(data, out if out.endswith(".pptx") else out + ".pptx")
    else:
        os.makedirs(out, exist_ok=True)
        export_csv(data, os.path.join(out, "题库.csv"))
        export_md(data, os.path.join(out, "试卷.md"))
        export_pptx(data, os.path.join(out, "问卷.pptx"))


if __name__ == "__main__":
    main()