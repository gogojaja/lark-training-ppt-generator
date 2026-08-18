# 流程图示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_flowchart(pres, title, nodes, edges, flow_type='standard'):
    """创建流程图"""
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 56, 100)

    # 创建流程节点
    start_x = Inches(2)
    start_y = Inches(2.5)
    node_width = Inches(1.5)
    node_height = Inches(0.6)

    # 绘制节点
    node_positions = []
    for i, node in enumerate(nodes):
        x = start_x + i * (node_width + Inches(0.5))
        y = start_y

        # 创建矩形形状
        shape = slide.shapes.add_shape(
            1,  # 矩形
            x, y, node_width, node_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(46, 117, 182)
        shape.line.color.rgb = RGBColor(31, 56, 100)

        # 设置圆角
        shape.shape_properties.preset = 13  # 圆角矩形

        # 添加文本
        text_frame = shape.text_frame
        text_frame.text = node
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        node_positions.append((x, y))

    # 绘制连接线
    for i in range(len(node_positions) - 1):
        x1, y1 = node_positions[i]
        x2, y2 = node_positions[i + 1]

        # 水平线
        line = slide.shapes.add_shape(
            1,  # 矩形
            x1 + node_width, y1 + node_height / 2 - 1,
            x2 - x1 - node_width, 2
        )
        line.line.color.rgb = RGBColor(31, 56, 100)
        line.line.width = Pt(2)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    nodes = ['开始', '签约申请', '材料审核', '合同签订', '完成']
    slide = create_flowchart(pres, '个人综合签约流程', nodes, [])
    pres.save('output/flowchart_example.pptx')
    print('流程图示例已生成: output/flowchart_example.pptx')
