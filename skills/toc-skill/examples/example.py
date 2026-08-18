# 目录页示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_toc_page(pres, title, items):
    """创建目录页"""
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 56, 100)

    # 目录项
    start_y = Inches(2.5)
    x_position = Inches(2)

    for i, item in enumerate(items):
        # 创建文本框
        text_box = slide.shapes.add_textbox(
            x_position, start_y + i * Inches(0.6), Inches(6), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = item

        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(24)
        text_para.font.color.rgb = RGBColor(51, 51, 51)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    items = [
        '个人综合签约概述 - 第1-2页',
        '签约流程详解 - 第3-8页',
        '注意事项 - 第9-12页',
        '常见问题 - 第13-15页'
    ]
    slide = create_toc_page(pres, '目录', items)
    pres.save('output/toc_example.pptx')
    print('目录页示例已生成: output/toc_example.pptx')
