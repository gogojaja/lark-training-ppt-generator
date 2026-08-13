# 表格示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_table(pres, title, columns, rows, table_type='standard'):
    """创建表格"""
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 56, 100)

    # 创建表格
    rows_count = len(rows) + 1
    cols_count = len(columns)

    table = slide.shapes.add_table(
        rows_count, cols_count, Inches(1), Inches(4), Inches(8), Inches(5)
    ).table

    # 设置列宽
    col_width = Inches(8) / cols_count
    for i in range(cols_count):
        table.columns[i].width = col_width

    # 表头
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(74, 111, 165)
        cell.text_frame.paragraphs[0].font.size = Pt(20)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 数据行
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(18)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 斑马纹效果
            if table_type == 'zebra' and i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 247, 250)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    columns = ['步骤', '操作', '时间', '负责人']
    rows = [
        ['1', '提交申请', '1-2天', '客户经理'],
        ['2', '审核材料', '2-3天', '风控专员'],
        ['3', '签订合同', '1天', '法务部门']
    ]
    slide = create_table(pres, '签约流程步骤对比', columns, rows, table_type='zebra')
    pres.save('output/table_example.pptx')
    print('表格示例已生成: output/table_example.pptx')
