# 封面页示例

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_cover_page(pres, title, subtitle, author, date):
    """创建封面页"""
    slide_layout = pres.slide_layouts[6]  # 空白布局
    slide = pres.slides.add_slide(slide_layout)

    # 设置背景色（深蓝渐变）
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 56, 100)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle

    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)

    # 作者
    author_box = slide.shapes.add_textbox(
        Inches(1), Inches(6), Inches(8), Inches(0.5)
    )
    author_frame = author_box.text_frame
    author_frame.text = author

    author_para = author_frame.paragraphs[0]
    author_para.alignment = PP_ALIGN.CENTER
    author_para.font.size = Pt(20)
    author_para.font.color.rgb = RGBColor(200, 200, 200)

    # 日期
    date_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5), Inches(8), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_frame.text = date

    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(20)
    date_para.font.color.rgb = RGBColor(200, 200, 200)

    return slide

# 使用示例
if __name__ == '__main__':
    pres = Presentation()
    slide = create_cover_page(
        pres,
        title='个人综合签约培训',
        subtitle='操作流程与注意事项',
        author='培训部',
        date='2026-08-13'
    )
    pres.save('output/cover_example.pptx')
    print('封面页示例已生成: output/cover_example.pptx')
